import os
from pptx import Presentation

# Directory where all the presentations are stored
presentations_dir = 'presentations'

# List all .pptx files in the presentations directory
presentation_files = [f for f in os.listdir(presentations_dir) if f.endswith('.pptx')]

# Iterate through each presentation file
for filename in presentation_files:
    # Construct the full file path
    file_path = os.path.join(presentations_dir, filename)
    
    # Open the PowerPoint file
    ppt = Presentation(file_path)
    
    print(f"Processing {filename}...")
    # Iterate through the slides and elements
    for slide_number, slide in enumerate(ppt.slides, start=1):
        print(f"  Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"    {shape.text}")
